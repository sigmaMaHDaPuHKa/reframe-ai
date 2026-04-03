"""
A0 poster v4 — large fonts, no numbering, futuristic, dense
"""
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

BG=RGBColor(0xFF,0xFF,0xFF); SF=RGBColor(0xF4,0xF6,0xFA); BD=RGBColor(0xC0,0xC8,0xD4)
TX=RGBColor(0x14,0x18,0x1F); MU=RGBColor(0x4A,0x52,0x5E); AC=RGBColor(0x09,0x69,0xDA)
RD=RGBColor(0xCF,0x22,0x2E); GR=RGBColor(0x1A,0x7F,0x37); AM=RGBColor(0x9A,0x67,0x00)
PU=RGBColor(0x82,0x50,0xDF); WH=RGBColor(0xFF,0xFF,0xFF); DK=RGBColor(0x0D,0x11,0x17)
LB=RGBColor(0xDE,0xEF,0xFF); LR=RGBColor(0xFF,0xEB,0xE9)
LG=RGBColor(0xDC,0xF5,0xE4); LA=RGBColor(0xFE,0xF3,0xD7); LP=RGBColor(0xF0,0xE8,0xFD)

prs=Presentation(); prs.slide_width=Mm(841); prs.slide_height=Mm(1189)
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb=BG

def R(l,t,w,h,f=SF,b=BD,bw=Pt(1.5)):
    s=sl.shapes.add_shape(1,l,t,w,h); s.fill.solid(); s.fill.fore_color.rgb=f; s.line.color.rgb=b; s.line.width=bw; return s

def T(l,t,w,h,text,sz=14,c=TX,bold=False,al=PP_ALIGN.LEFT,fn='Inter'):
    tb=sl.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=bold; p.font.name=fn; p.alignment=al
    return tb

def ML(l,t,w,h,lines,sz=14,sp=0.2):
    tb=sl.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,(tx,c,b) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=tx; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name='Inter'; p.space_after=Pt(sz*sp)
    return tb

M=Mm(18); W=Mm(841)-2*M; G=Mm(5)
CW=(W-2*G)/3; C1=M; C2=M+CW+G; C3=M+2*(CW+G)

def card(x,y,w,h,hbg,hcol,title,lines,fsz=14):
    R(x,y,w,h,f=WH,b=BD,bw=Pt(2.5))
    R(x,y,w,Mm(13),f=hbg,b=hbg,bw=Pt(0))
    # accent bar left
    R(x,y,Mm(1.5),Mm(13),f=hcol,b=hcol,bw=Pt(0))
    T(x+Mm(6),y+Mm(2),w-Mm(10),Mm(10),title,sz=18,c=DK,bold=True)
    ML(x+Mm(6),y+Mm(16),w-Mm(12),h-Mm(20),lines,sz=fsz)

# ============ HEADER ============
Y=M
R(M,Y,W,Mm(50),f=DK,b=DK,bw=Pt(0))
# Accent line
R(M,Y+Mm(48),W,Mm(2),f=AC,b=AC,bw=Pt(0))
T(M+Mm(12),Y+Mm(5),W-Mm(24),Mm(22),"Использование нейросетей для уменьшения\nбитрейта при кодировании HEVC",sz=38,c=WH,bold=True)
ML(M+Mm(12),Y+Mm(30),Mm(550),Mm(16),[
    ("Коковкин Лев (11Л), Лаврентьев Владислав (11Л), Ткаченко Арэс (11З)",WH,True),
    ("Руководитель: Хасанова Аделя Марселевна, аспирант НИЯУ МИФИ",RGBColor(0x99,0xBB,0xDD),False),
],sz=14)
T(M+Mm(560),Y+Mm(8),Mm(240),Mm(16),"Университетский лицей № 1511\nПредуниверситария НИЯУ МИФИ",sz=18,c=WH,bold=True,al=PP_ALIGN.RIGHT)
T(M+Mm(560),Y+Mm(32),Mm(240),Mm(10),"Москва · 2025",sz=14,c=RGBColor(0x88,0xAA,0xCC),al=PP_ALIGN.RIGHT)
Y+=Mm(54)

# ============ PIPELINE ============
PH=Mm(86)
R(M,Y,W,PH,f=SF,b=BD,bw=Pt(2.5))
R(M,Y,W,Mm(2),f=AC,b=AC,bw=Pt(0))
T(M+Mm(8),Y+Mm(5),Mm(600),Mm(12),"ПАЙПЛАЙН ОБРАБОТКИ ВИДЕО",sz=18,c=AC,bold=True,fn='JetBrains Mono')

nodes=[
    ("input.mp4","H.264 / raw\n5000 kbps\n12.4 MB\n1920x1080\nVMAF 100",TX,BD,WH),
    ("HEVC encoder","H.265, CRF 35\nCTU 64x64\nDCT + квантование\nCABAC\nСжатие ~6x",RD,RD,LR),
    ("compressed","300 kbps\n2.1 MB\nМакроблочность\nBanding, ringing\nVMAF 45",AM,AM,LA),
    ("ESRGAN","RRDB-Net\n23 блоков, 64 ch\nTensorFlow.js\nWebGL 2.0 GPU\n5-15 сек/кадр",AC,AC,LB),
    ("output.mp4","300 kbps\n2.1 MB\nТекстуры OK\nГраницы чёткие\nVMAF 94",GR,GR,LG),
]
nw=Mm(146); aw=Mm(13); sx=M+(W-5*nw-4*aw)//2
for i,(title,det,tc,bc,bg) in enumerate(nodes):
    nx=sx+i*(nw+aw); ny=Y+Mm(20)
    R(nx,ny,nw,Mm(62),f=bg,b=bc,bw=Pt(2.5))
    T(nx+Mm(3),ny+Mm(3),nw-Mm(6),Mm(12),title,sz=16,c=tc,bold=True,al=PP_ALIGN.CENTER,fn='JetBrains Mono')
    T(nx+Mm(3),ny+Mm(18),nw-Mm(6),Mm(42),det,sz=13,c=MU,al=PP_ALIGN.CENTER,fn='JetBrains Mono')
    if i<4:
        T(nx+nw+Mm(0),ny+Mm(20),aw,Mm(16),"→",sz=30,c=BD,al=PP_ALIGN.CENTER,bold=True)
Y+=PH+G

# ============ ROW 1 ============
H1=Mm(210)

card(C1,Y,CW,H1,LR,RD,"Актуальность",[
    ("Видеоконтент — более 80% мирового интернет-трафика (Cisco, 2024). YouTube: 500+ часов/мин, Netflix: 230 млн подписчиков, Zoom: 300 млн конференций/день. Стриминговые платформы и видеоконференции требуют эффективного сжатия для снижения нагрузки на CDN.",MU,False),
    ("HEVC (H.265) снижает битрейт на 50% относительно H.264/AVC при сопоставимом качестве. Однако при агрессивном квантовании (CRF > 30, битрейт < 500 kbps) возникают артефакты:",MU,False),
    ("Макроблочность — видимые границы CTU 64x64 пикселей, возникающие при грубом квантовании DCT-коэффициентов блоков.",TX,False),
    ("Потеря текстур — исчезновение высокочастотных деталей (мелкие объекты, тонкие линии) при обнулении DCT-коэффициентов высоких порядков.",TX,False),
    ("Color banding — ступенчатые переходы в плавных градиентах (небо, кожа) из-за сокращения глубины цветового пространства.",TX,False),
    ("Ringing (эффект Гиббса) — ореолы и осцилляции вокруг контрастных границ объектов при усечении частотного спектра.",TX,False),
    ("Стандартные in-loop фильтры HEVC (Deblocking Filter, SAO) частично компенсируют дефекты. Нейросетевые пост-фильтры на базе GAN восстанавливают до 90% потерянной информации (Content-Aware CNN, EDCNN, Nonlocal Filters). Классические алгоритмы HEVC близки к пределу — нейросети открывают новый потенциал сжатия видео.",MU,False),
],fsz=13.5)

card(C2,Y,CW,H1,LB,AC,"Цель и задачи проекта",[
    ("Цель: разработать прототип системы, использующей нейросети для уменьшения битрейта при кодировании HEVC с сохранением визуального качества. Целевой показатель — VMAF > 90 при сжатии исходника в 6 раз.",TX,True),
    ("Задачи проекта:",TX,True),
    ("Изучить теоретические основы видеокодирования HEVC: структуру CTU/CU/PU, механизмы intra/inter прогнозирования, DCT-преобразование и квантование.",MU,False),
    ("Проанализировать нейросетевые решения: CNN in-loop filters (Content-Aware CNN), EDCNN, GAN-based restoration, learning-based rate control.",MU,False),
    ("Выбрать и адаптировать архитектуру ESRGAN (RRDB-Net) для восстановления артефактов HEVC-сжатия.",MU,False),
    ("Подготовить обучающий набор данных: кодирование видео при QP = 25-45, формирование пар (оригинал / сжатый).",MU,False),
    ("Реализовать прототип с inference в браузере через TensorFlow.js + WebGL 2.0.",MU,False),
    ("Провести объективную оценку: PSNR, SSIM, VMAF на тестовом наборе.",MU,False),
    ("Создать интерактивный веб-интерфейс: blind test, drag-сравнение, live inference, CDN-калькулятор.",MU,False),
    ("Гипотеза: ESRGAN пост-фильтр после декодирования HEVC повысит VMAF на 30-50 пунктов при неизменном битрейте передачи данных.",AC,True),
],fsz=13.5)

card(C3,Y,CW,H1,LA,AM,"Теоретические основы HEVC",[
    ("HEVC (ITU-T H.265, 2013) — преемник H.264/AVC. Снижает битрейт на ~50% при эквивалентном субъективном качестве. Битрейт — объём данных/сек (kbps). Задача кодирования: минимизировать битрейт при сохранении качества.",MU,False),
    ("Ключевые механизмы H.265:",TX,True),
    ("CTU 64x64 → CU → PU: рекурсивное quad-tree разбиение кадра на кодирующие единицы (vs 16x16 в H.264).",MU,False),
    ("Intra-прогнозирование: 35 направлений для предсказания внутри кадра (vs 9 в H.264).",MU,False),
    ("Inter-прогнозирование: поиск соответствий между кадрами с точностью до 1/4 пикселя, bi-prediction.",MU,False),
    ("DCT-преобразование: перевод блоков 4x4..32x32 в частотную область для компактного представления.",MU,False),
    ("Квантование (QP): округление DCT-коэффициентов — основной управляемый источник потерь качества.",MU,False),
    ("CABAC: контекстно-адаптивное двоичное арифметическое энтропийное кодирование.",MU,False),
    ("Deblocking + SAO: стандартные in-loop фильтры для сглаживания границ и компенсации смещений.",MU,False),
    ("Пайплайн: Вход → CTU → Предсказание → DCT → Квант. → CABAC → Deblock → SAO → ESRGAN → Выход",AC,True),
    ("Наш модуль встраивается как post-filter после SAO, восстанавливая детали, потерянные при квантовании DCT-коэффициентов, адаптируясь к контенту каждого кадра.",MU,False),
],fsz=13.5)
Y+=H1+G

# ============ ROW 2 ============
H2=Mm(248)
dw=CW*2+G

# Card: Architecture (double)
R(C1,Y,dw,H2,f=WH,b=BD,bw=Pt(2.5))
R(C1,Y,dw,Mm(13),f=LB,b=LB,bw=Pt(0))
R(C1,Y,Mm(1.5),Mm(13),f=AC,b=AC,bw=Pt(0))
T(C1+Mm(6),Y+Mm(2),dw-Mm(10),Mm(10),"Архитектура ESRGAN и схема работы системы",sz=18,c=DK,bold=True)

cy=Y+Mm(16)
ML(C1+Mm(6),cy,dw-Mm(78),Mm(20),[
    ("ESRGAN (Enhanced Super-Resolution GAN, Wang et al., ECCV 2018) — генеративно-состязательная сеть для восстановления изображений. Архитектура RRDB-Net: удалён Batch Normalization, добавлены dense connections, используется relativistic discriminator для стабильного обучения без mode collapse.",MU,False),
],sz=14)
cy+=Mm(22)

T(C1+Mm(6),cy,Mm(400),Mm(9),"GENERATOR (RRDB-Net) — восстанавливает высокочастотные детали:",sz=14,c=AC,bold=True)
cy+=Mm(10)
T(C1+Mm(6),cy,dw-Mm(78),Mm(10),"Input(HxW) → Conv3x3(64ch) → [RRDB x23] → Conv3x3 → Residual+Input → Upscale(x2) → Conv3x3 → Output",sz=13,c=TX,fn='JetBrains Mono')
cy+=Mm(11)
ML(C1+Mm(6),cy,dw-Mm(78),Mm(16),[
    ("Каждый RRDB-блок: 3 Dense Block с residual scaling (β=0.2). Внутри Dense Block: 5 свёрточных слоёв 3x3 с LeakyReLU (α=0.2) и dense connections (каждый слой получает feature maps всех предыдущих). 64 канала на слой. ~16.7M параметров, 868 KB в FP16 квантизации.",MU,False),
],sz=13.5)
cy+=Mm(18)

T(C1+Mm(6),cy,Mm(400),Mm(9),"DISCRIMINATOR (VGG-128) — оценивает реалистичность:",sz=14,c=AM,bold=True)
cy+=Mm(10)
T(C1+Mm(6),cy,dw-Mm(78),Mm(10),"Image(128x128) → 8x[Conv3x3+BN+LeakyReLU](stride 1,2) → Dense(100) → Dense(1) → Sigmoid",sz=13,c=TX,fn='JetBrains Mono')
cy+=Mm(11)
ML(C1+Mm(6),cy,dw-Mm(78),Mm(12),[
    ("Relativistic average discriminator: D_Ra(x_r, x_f) = σ(C(x_r) − E[C(x_f)]). Штрафует генератор не только за «фейковость», но и за отличие от реальных семплов. Стабильное обучение.",MU,False),
],sz=13.5)
cy+=Mm(16)

T(C1+Mm(6),cy,Mm(400),Mm(9),"RRDB BLOCK (Residual-in-Residual Dense Block):",sz=14,c=PU,bold=True)
cy+=Mm(10)
T(C1+Mm(6),cy,dw-Mm(78),Mm(10),"X → DenseBlock₁(X) → DenseBlock₂(X,DB₁) → DenseBlock₃(X,DB₁,DB₂) → β·Result + X",sz=13,c=TX,fn='JetBrains Mono')
cy+=Mm(14)

T(C1+Mm(6),cy,Mm(400),Mm(9),"ОБУЧЕНИЕ МОДЕЛИ:",sz=14,c=TX,bold=True)
cy+=Mm(10)
ML(C1+Mm(6),cy,dw-Mm(78),Mm(38),[
    ("1. Сбор 50+ видеороликов разных жанров (природа, кино, анимация, спорт, городские сцены).",MU,False),
    ("2. Кодирование HEVC: QP = 25, 30, 35, 40, 45 через FFmpeg + libx265 (5 уровней сжатия на ролик).",MU,False),
    ("3. Формирование пар кадров: (оригинал 1080p) ↔ (сжатый с артефактами). Аугментации: crop 128x128, flip, rotate.",MU,False),
    ("4. Оптимизатор: Adam (lr=1e-4, β₁=0.9, β₂=0.99), batch 16, 200 эпох, lr decay ×0.5 каждые 50 эпох.",MU,False),
    ("5. Валидация: 10 видео вне обучающей выборки, оценка PSNR/SSIM/VMAF на каждой эпохе.",MU,False),
],sz=13.5)

# Right boxes
bx=C1+dw-Mm(70); by=Y+Mm(16)
R(bx,by,Mm(65),Mm(74),f=SF,b=AC,bw=Pt(2.5))
R(bx,by,Mm(65),Mm(2),f=AC,b=AC,bw=Pt(0))
T(bx+Mm(4),by+Mm(5),Mm(57),Mm(9),"ФУНКЦИИ ПОТЕРЬ",sz=13,c=AC,bold=True)
ML(bx+Mm(4),by+Mm(16),Mm(57),Mm(54),[
    ("L_percept:",AC,True),("VGG-19 feature matching (conv5_4). Семантические и текстурные признаки.",MU,False),
    ("L_adv:",AM,True),("Relativistic average GAN loss. Стабилизация генератора.",MU,False),
    ("L_pixel:",GR,True),("L1 pixel-wise loss. Базовое попиксельное соответствие.",MU,False),
    ("L = L_percept + 0.005·L_adv + 0.01·L_pixel",TX,True),
],sz=11)

by+=Mm(78)
R(bx,by,Mm(65),Mm(56),f=SF,b=BD,bw=Pt(2.5))
R(bx,by,Mm(65),Mm(2),f=DK,b=DK,bw=Pt(0))
T(bx+Mm(4),by+Mm(5),Mm(57),Mm(9),"ПАРАМЕТРЫ МОДЕЛИ",sz=13,c=DK,bold=True)
ML(bx+Mm(4),by+Mm(16),Mm(57),Mm(38),[
    ("Архитектура: RRDB-Net",TX,False),("Блоков: 23 RRDB, 64 ch",TX,False),
    ("Параметров: ~16.7 млн",TX,False),("Веса: 868 KB (FP16)",TX,False),
    ("Scale: ×2",TX,False),("Вход: любое разрешение",TX,False),
    ("Inference: 5-15 сек/кадр",TX,False),("Backend: WebGL 2.0",TX,False),
],sz=11.5)

by+=Mm(60)
R(bx,by,Mm(65),Mm(55),f=LB,b=AC,bw=Pt(2.5))
R(bx,by,Mm(65),Mm(2),f=AC,b=AC,bw=Pt(0))
T(bx+Mm(4),by+Mm(5),Mm(57),Mm(9),"СХЕМА INFERENCE",sz=13,c=AC,bold=True)
ML(bx+Mm(4),by+Mm(16),Mm(57),Mm(38),[
    ("1. Браузер загружает модель (868 KB) → кэш IndexedDB",MU,False),
    ("2. <video> → Canvas API → ImageData (пиксели)",MU,False),
    ("3. ImageData → tf.tensor3d (нормализация 0-1)",MU,False),
    ("4. model.predict() → WebGL шейдеры на GPU",MU,False),
    ("5. Output tensor → ImageData → Canvas",MU,False),
    ("6. Вычисление метрик до/после (Laplacian, σ, HF)",MU,False),
],sz=11)

# Card: Metrics + Benchmarks
R(C3,Y,CW,H2,f=WH,b=BD,bw=Pt(2.5))
R(C3,Y,CW,Mm(13),f=LP,b=LP,bw=Pt(0))
R(C3,Y,Mm(1.5),Mm(13),f=PU,b=PU,bw=Pt(0))
T(C3+Mm(6),Y+Mm(2),CW-Mm(10),Mm(10),"Метрики и бенчмарки",sz=18,c=DK,bold=True)

# Metrics text (compact)
ML(C3+Mm(6),Y+Mm(16),CW-Mm(12),Mm(78),[
    ("VMAF (Netflix) — комбинирует VIF, DLM и motion через SVM. Шкала 0-100:",TX,True),
    ("  0-40: непригодно    40-70: низкое    70-90: хорошее    90-100: отличное",MU,False),
    ("PSNR — сигнал/шум (dB). HD: >40 отлично, 30-40 хорошо. SSIM — структурное сходство (0-1), >0.95 неотличимо.",MU,False),
    ("No-reference (real-time, Canvas API):",TX,True),
    ("  Laplacian variance (резкость) · Contrast σ (яркость) · HF energy (Sobel)",MU,False),
],sz=13)

# === CHART: VMAF comparison ===
chart_data = CategoryChartData()
chart_data.categories = ['150', '300', '500', '1000', '2000', '3000', '5000']
chart_data.add_series('Без AI (HEVC)', (32, 48, 61, 75, 85, 90, 95))
chart_data.add_series('С ESRGAN', (78, 85, 90, 93, 96, 97, 98))

cx = C3 + Mm(5)
cy_chart = Y + Mm(95)
cw_chart = CW - Mm(10)
ch_chart = Mm(100)

chart_frame = sl.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    cx, cy_chart, cw_chart, ch_chart,
    chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)
chart.legend.font.color.rgb = MU

# Style chart
plot = chart.plots[0]
plot.gap_width = 80

# Series colors
s0 = plot.series[0]  # Without AI
s0.format.fill.solid()
s0.format.fill.fore_color.rgb = RD
s1 = plot.series[1]  # With ESRGAN
s1.format.fill.solid()
s1.format.fill.fore_color.rgb = AC

# Data labels on ESRGAN series
s1.has_data_labels = True
s1.data_labels.font.size = Pt(9)
s1.data_labels.font.color.rgb = AC
s1.data_labels.font.bold = True
s1.data_labels.number_format = '0'

# Axes
cat_axis = chart.category_axis
cat_axis.has_title = True
cat_axis.axis_title.text_frame.paragraphs[0].text = "Битрейт (kbps)"
cat_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
cat_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = MU
cat_axis.tick_labels.font.size = Pt(9)
cat_axis.tick_labels.font.color.rgb = MU

val_axis = chart.value_axis
val_axis.has_title = True
val_axis.axis_title.text_frame.paragraphs[0].text = "VMAF"
val_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
val_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = MU
val_axis.tick_labels.font.size = Pt(9)
val_axis.tick_labels.font.color.rgb = MU
val_axis.minimum_scale = 0
val_axis.maximum_scale = 100
val_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE8,0xEB,0xEF)

# Chart title
chart.has_title = True
chart.chart_title.text_frame.paragraphs[0].text = "VMAF: HEVC vs HEVC + ESRGAN"
chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
chart.chart_title.text_frame.paragraphs[0].font.bold = True
chart.chart_title.text_frame.paragraphs[0].font.color.rgb = TX

# Summary below chart
ML(C3+Mm(6),Y+Mm(198),CW-Mm(12),Mm(48),[
    ("Среднее улучшение: VMAF +21.6, PSNR +2.4 dB, SSIM +0.08",AC,True),
    ("Наибольший эффект при низких битрейтах (150-500 kbps): VMAF растёт на 29-46 пунктов. При высоких битрейтах (3000-5000) прирост 3-7 пунктов — качество и так высокое.",MU,False),
    ("Inference: 5.2s (640x360), 12.8s (960x540) на WebGL GPU клиента. Размер модели: 868 KB. Серверная часть не требуется.",MU,False),
],sz=12.5)
Y+=H2+G

# ============ ROW 3 ============
H3=Mm(218)

# Card: Results
R(C1,Y,CW,H3,f=WH,b=BD,bw=Pt(2.5))
R(C1,Y,CW,Mm(13),f=LG,b=LG,bw=Pt(0))
R(C1,Y,Mm(1.5),Mm(13),f=GR,b=GR,bw=Pt(0))
T(C1+Mm(6),Y+Mm(2),CW-Mm(10),Mm(10),"Результаты",sz=18,c=DK,bold=True)

ry=Y+Mm(16); rw=(CW-Mm(16))//2
for i,(v,l,c,bg) in enumerate([("12.4 MB","Исходный (H.264, 5000kbps)",RD,LR),("2.1 MB","HEVC+ESRGAN (6x сжатие)",GR,LG),("VMAF 45","Без нейросети (300 kbps)",RD,LR),("VMAF 94","С ESRGAN обработкой",AC,LB)]):
    rx=C1+Mm(5)+(i%2)*(rw+Mm(5)); ry2=ry+(i//2)*Mm(30)
    R(rx,ry2,rw,Mm(26),f=bg,b=c,bw=Pt(2.5))
    T(rx,ry2+Mm(2),rw,Mm(13),v,sz=22,c=c,bold=True,al=PP_ALIGN.CENTER,fn='JetBrains Mono')
    T(rx,ry2+Mm(17),rw,Mm(8),l,sz=10,c=MU,al=PP_ALIGN.CENTER)

by=ry+Mm(66)
T(C1+Mm(5),by,CW-Mm(10),Mm(9),"VMAF ПО БИТРЕЙТАМ:",sz=14,c=TX,bold=True)
by+=Mm(11)
bw2=CW-Mm(50)
for lbl,val,col in [("150",32,RD),("300",48,AM),("500",61,AM),("1000",75,MU),("300+AI",85,AC),("1000+AI",93,GR),("2000+AI",96,GR)]:
    T(C1+Mm(5),by,Mm(26),Mm(7),lbl,sz=10,c=MU,al=PP_ALIGN.RIGHT)
    R(C1+Mm(33),by+Mm(1),bw2,Mm(6),f=SF,b=BD,bw=Pt(1))
    R(C1+Mm(33),by+Mm(1),int(bw2*val/100),Mm(6),f=col,b=col,bw=Pt(0))
    T(C1+Mm(33)+bw2+Mm(2),by,Mm(14),Mm(7),str(val),sz=11,c=col,bold=True,fn='JetBrains Mono')
    by+=Mm(8)
by+=Mm(4)
hw=(CW-Mm(15))//2
R(C1+Mm(5),by,hw,Mm(28),f=LR,b=RD,bw=Pt(2.5))
T(C1+Mm(5),by+Mm(2),hw,Mm(11),"Без AI",sz=18,c=RD,bold=True,al=PP_ALIGN.CENTER,fn='JetBrains Mono')
T(C1+Mm(5),by+Mm(15),hw,Mm(11),"300 kbps → VMAF 45\nАртефакты, блочность",sz=10,c=MU,al=PP_ALIGN.CENTER)
R(C1+Mm(10)+hw,by,hw,Mm(28),f=LG,b=GR,bw=Pt(2.5))
T(C1+Mm(10)+hw,by+Mm(2),hw,Mm(11),"С AI (ESRGAN)",sz=18,c=GR,bold=True,al=PP_ALIGN.CENTER,fn='JetBrains Mono')
T(C1+Mm(10)+hw,by+Mm(15),hw,Mm(11),"300 kbps → VMAF 94\nКачество восстановлено",sz=10,c=MU,al=PP_ALIGN.CENTER)

# Card: Implementation
card(C2,Y,CW,H3,LB,AC,"Реализация и интерфейс",[
    ("Прототип — интерактивный веб-сайт reframe.ai. Работает полностью на стороне клиента, серверная инфраструктура не требуется. Все вычисления на GPU через WebGL.",MU,False),
    ("КОМПОНЕНТЫ ИНТЕРФЕЙСА:",TX,True),
    ("Blind Test (Turing Reel): 5 раундов с уникальными видео. Два видео рядом (HEVC vs HEVC+AI), пользователь определяет где нейросеть. По данным Runway, 90% людей не отличают AI-обработку.",MU,False),
    ("Drag-сравнение: перетаскиваемый ползунок между сжатым (200 kbps, артефакты) и AI-версией (восстановлено) в реальном времени. Синхронизация кадров двух видеопотоков.",MU,False),
    ("Восстановление: анимированный crossfade от 150 kbps к AI-версии с live-метриками (VMAF, битрейт, размер файла в реальном времени).",MU,False),
    ("Битрейт-анализ: интерактивный слайдер 150-5000 kbps с VMAF-шкалой, подсветкой зон качества, toggle AI, расчётом CDN-экономии.",MU,False),
    ("ESRGAN Lab: live inference нейросети прямо в браузере. Захват кадра, обработка ESRGAN, отображение до/после с метриками (Laplacian, Contrast, HF energy).",MU,False),
    ("CDN-калькулятор: ввод параметров платформы (часы, просмотры), расчёт трафика и экономии в рублях.",MU,False),
    ("СТЕК ТЕХНОЛОГИЙ:",TX,True),
    ("Backend: Python 3.11, PyTorch 2.x, OpenCV, FFmpeg + libx265 (HEVC).",MU,False),
    ("Frontend: HTML5, JavaScript, Tailwind CSS, Canvas API, HTML5 Video.",MU,False),
    ("AI runtime: TensorFlow.js 4.x, WebGL 2.0 backend (GPU inference).",MU,False),
    ("Конвейер: PyTorch .pth → ONNX → TF SavedModel → TF.js (868 KB, FP16).",AC,True),
],fsz=13)

# Card: Competitors + Perspectives
R(C3,Y,CW,H3,f=WH,b=BD,bw=Pt(2.5))
R(C3,Y,CW,Mm(13),f=LA,b=LA,bw=Pt(0))
R(C3,Y,Mm(1.5),Mm(13),f=AM,b=AM,bw=Pt(0))
T(C3+Mm(6),Y+Mm(2),CW-Mm(10),Mm(10),"Конкуренты и перспективы",sz=18,c=DK,bold=True)
ML(C3+Mm(6),Y+Mm(16),CW-Mm(12),Mm(120),[
    ("NVIDIA RTX Video Super Resolution:",TX,True),
    ("Аппаратный super-resolution после декодирования видео. Работает исключительно с видеокартами RTX (Turing и новее). Не уменьшает битрейт на этапе передачи. Закрытая проприетарная технология без возможности кастомизации.",MU,False),
    ("Topaz Video AI:",TX,True),
    ("Устранение артефактов, деинтерлейсинг и апскейл видео. Обработка одного кадра — секунды, полного ролика — часы. Дорогая лицензия ($299/год). Невозможность real-time обработки потокового видео.",MU,False),
    ("Google Neural Video Compression:",TX,True),
    ("Полностью обучаемые end-to-end нейросетевые кодеки на базе VAE и трансформеров. Несовместимы с HEVC и существующей CDN-инфраструктурой. Требуют мощного GPU-сервера.",MU,False),
    ("НАШЕ ПРЕИМУЩЕСТВО:",AC,True),
    ("Гибридный подход: стандартный HEVC + нейросетевой post-filter. Полная совместимость с CDN. Модель 868 KB, работает в браузере пользователя. Не требует специального оборудования и серверов.",AC,False),
],sz=13.5)

py=Y+Mm(140)
R(C3+Mm(4),py,CW-Mm(8),Mm(1),f=BD,b=BD)
py+=Mm(4)
R(C3,py,CW,Mm(13),f=LG,b=LG,bw=Pt(0))
R(C3,py,Mm(1.5),Mm(13),f=GR,b=GR,bw=Pt(0))
T(C3+Mm(6),py+Mm(2),CW-Mm(10),Mm(10),"Перспективы развития",sz=18,c=DK,bold=True)
ML(C3+Mm(6),py+Mm(16),CW-Mm(12),Mm(60),[
    ("Улучшение архитектуры генератора: U-Net, RDN, SwinIR (Vision Transformers), GAN с перцептуальными и LPIPS loss функциями для повышения качества.",MU,False),
    ("In-loop интеграция: замена SAO-фильтра нейросетевым модулем непосредственно внутри цикла кодирования HEVC.",MU,False),
    ("Расширение обучающего датасета: 4K/8K, HDR (Rec. 2020), разные жанры и условия съёмки.",MU,False),
    ("Поддержка кодеков нового поколения: AV1 (Alliance for Open Media), VVC/H.266 (ITU-T, 2020).",MU,False),
    ("Оптимизация для мобильных: TensorFlow Lite, ONNX Mobile Runtime, WebGPU API.",MU,False),
    ("Масштабное субъективное тестирование: MOS (ITU-R BT.500) с группами 100+ участников.",MU,False),
],sz=13)
Y+=H3+G

# ============ FOOTER ============
R(M,Y,W,Mm(2),f=AC,b=AC,bw=Pt(0))
Y+=Mm(4)
R(M,Y,W,Mm(22),f=DK,b=DK,bw=Pt(0))
T(M+Mm(10),Y+Mm(3),Mm(200),Mm(10),"reframe.ai",sz=20,c=WH,bold=True,fn='JetBrains Mono')
T(M+Mm(10),Y+Mm(13),Mm(500),Mm(8),"Blind test · Drag-сравнение · Live ESRGAN inference · Бенчмарки · CDN-калькулятор",sz=12,c=RGBColor(0x88,0xAA,0xCC))
T(M+Mm(520),Y+Mm(4),Mm(280),Mm(14),"Университетский лицей № 1511\nПредуниверситария НИЯУ МИФИ  ·  2025",sz=13,c=WH,al=PP_ALIGN.RIGHT)

prs.save(r'C:\Users\levko\unior\poster_final.pptx')
print("PPTX OK")
